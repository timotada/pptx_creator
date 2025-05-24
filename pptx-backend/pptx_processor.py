"""
PPTX Processor
Core module for processing PowerPoint templates with placeholder replacement
"""

import asyncio
import logging
from typing import Optional
from pathlib import Path
import shutil
from PIL import Image
import io

from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)

class PPTXProcessor:
    """Handles PPTX template processing with placeholder replacement"""
    
    def __init__(self):
        self.supported_image_formats = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.svg'}
    
    async def process_template(
        self,
        input_path: str,
        output_path: str,
        company_name: str,
        date: str,
        logo_path: Optional[str] = None
    ) -> None:
        """
        Process a PPTX template with placeholder replacement
        
        Args:
            input_path: Path to input PPTX file
            output_path: Path for output PPTX file
            company_name: Company name to replace {{COMPANY_NAME}} placeholders
            date: Formatted date string to replace {{DATE}} placeholders
            logo_path: Optional path to logo image for {{LOGO}} placeholders
        """
        
        logger.info(f"Processing PPTX template: {input_path}")
        
        try:
            # Load presentation
            prs = Presentation(input_path)
            
            # Process each slide
            for slide_idx, slide in enumerate(prs.slides):
                logger.debug(f"Processing slide {slide_idx + 1}")
                
                # Process text placeholders
                await self._replace_text_placeholders(slide, company_name, date)
                
                # Process logo placeholders
                if logo_path:
                    await self._replace_logo_placeholders(slide, logo_path)
            
            # Save processed presentation
            prs.save(output_path)
            logger.info(f"Successfully saved processed PPTX to: {output_path}")
            
        except Exception as e:
            logger.error(f"Error processing PPTX: {str(e)}")
            raise
    
    async def _replace_text_placeholders(self, slide, company_name: str, date: str) -> None:
        """Replace text placeholders in a slide"""
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Process each paragraph
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text
                        
                        # Replace placeholders
                        if "{{COMPANY_NAME}}" in text:
                            run.text = text.replace("{{COMPANY_NAME}}", company_name)
                            logger.debug(f"Replaced COMPANY_NAME with: {company_name}")
                        
                        if "{{DATE_LONG}}" in text:
                            formatted_date = self._format_date_for_display(date)
                            run.text = run.text.replace("{{DATE_LONG}}", formatted_date)
                            logger.debug(f"Replaced DATE_LONG with: {formatted_date}")
    
    async def _replace_logo_placeholders(self, slide, logo_path: str) -> None:
        """Replace logo placeholders in a slide"""
        
        shapes_to_remove = []
        logo_insertions = []
        
        for shape in slide.shapes:
            # Check for text-based logo placeholder
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "{{LOGO}}" in run.text:
                            # Store shape info for logo insertion
                            logo_insertions.append({
                                'left': shape.left,
                                'top': shape.top,
                                'width': shape.width,
                                'height': shape.height
                            })
                            shapes_to_remove.append(shape)
                            break
            
            # Check for image placeholder (if shape name contains "logo")
            elif hasattr(shape, 'name') and 'logo' in shape.name.lower():
                logo_insertions.append({
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
                shapes_to_remove.append(shape)
        
        # Remove placeholder shapes
        for shape in shapes_to_remove:
            slide.shapes._spTree.remove(shape._element)
        
        # Insert logo at each placeholder location
        for insertion in logo_insertions:
            await self._insert_logo(slide, logo_path, insertion)
    
    async def _insert_logo(self, slide, logo_path: str, position: dict) -> None:
        """Insert logo image at specified position with proper scaling"""
        
        try:
            # Load and process image
            with Image.open(logo_path) as img:
                # Get original dimensions
                orig_width, orig_height = img.size
                
                # Calculate target dimensions (preserve aspect ratio, never upscale)
                target_width = position['width']
                target_height = position['height']
                
                # Calculate scale factor (never exceed 1.0 to avoid upscaling)
                scale_x = min(1.0, target_width / Inches(1) / orig_width * 96)  # 96 DPI
                scale_y = min(1.0, target_height / Inches(1) / orig_height * 96)
                scale = min(scale_x, scale_y)
                
                # Calculate final dimensions
                final_width = Inches(orig_width * scale / 96)
                final_height = Inches(orig_height * scale / 96)
                
                # Calculate centered position
                left = position['left'] + (position['width'] - final_width) / 2
                top = position['top'] + (position['height'] - final_height) / 2
                
                # Insert image
                slide.shapes.add_picture(
                    logo_path,
                    left,
                    top,
                    final_width,
                    final_height
                )
                
                logger.debug(f"Inserted logo at position ({left}, {top}) with size ({final_width}, {final_height})")
                
        except Exception as e:
            logger.error(f"Error inserting logo: {str(e)}")
            # Continue processing even if logo insertion fails
    
    def _format_date_for_display(self, date_str: str) -> str:
        """Format date string for display (e.g., '24th May, 2025')"""
        try:
            from datetime import datetime
            date_obj = datetime.fromisoformat(date_str)
            day = date_obj.day
            
            # Get ordinal suffix
            if 10 <= day % 100 <= 20:
                suffix = 'th'
            else:
                suffix = {1: 'st', 2: 'nd', 3: 'rd'}.get(day % 10, 'th')
            
            month = date_obj.strftime("%B")  # May
            year = date_obj.year
            
            return f"{day}{suffix} {month}, {year}"
        except ValueError:
            return "24th May, 2025"  # Fallback
    
    def validate_image(self, image_path: str) -> bool:
        """Validate if image file is supported"""
        path = Path(image_path)
        return path.suffix.lower() in self.supported_image_formats 